"""
Generation MCP Server.

Exposes two tools:
  make_pdf(title, sections)  -> {path: str}
  make_pptx(title, slides)   -> {path: str}

This server renders already-structured content into files.
It does NOT summarize, transcribe, or analyze — that is the orchestrator's job.
Inputs are plain dicts of strings; outputs are file paths on disk.

Input shapes:
  make_pdf(
    title:    str,
    sections: [{"heading": str, "bullets": [str]}]
  )

  make_pptx(
    title:  str,
    slides: [{"heading": str, "bullets": [str]}]
  )

Outputs are written to the /outputs directory at the repo root.
File names are timestamped so repeated calls don't overwrite each other.

Run as a subprocess; communicates via stdio JSON-RPC (see mcp_common/server.py).
"""

import sys
from pathlib import Path
from datetime import datetime

sys.path.insert(0, str(Path(__file__).parent.parent.parent))

from backend.mcp_common.server import MCPServer

_OUTPUTS_DIR = Path(__file__).parent.parent.parent / "outputs"


def _ensure_outputs_dir() -> Path:
    _OUTPUTS_DIR.mkdir(exist_ok=True)
    return _OUTPUTS_DIR


def _timestamped_path(stem: str, suffix: str) -> Path:
    """Return outputs/<stem>_YYYYMMDD_HHMMSS.<suffix>."""
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    return _ensure_outputs_dir() / f"{stem}_{ts}{suffix}"


# Tool: make_pdf


def make_pdf(title: str, sections: list) -> dict:
    """
    Render a PDF from structured sections using fpdf2.

    Args:
      title     Document title, shown on the first page.
      sections  List of {"heading": str, "bullets": [str]}.
                Each section becomes a heading + bulleted list.

    Returns:
      path  Absolute path to the generated .pdf file.
    """
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    # Title page header
    pdf.set_font("Helvetica", style="B", size=20)
    pdf.set_fill_color(30, 30, 30)
    pdf.set_text_color(255, 255, 255)
    # Full-width title bar
    pdf.cell(0, 14, _safe(title), new_x="LMARGIN", new_y="NEXT", fill=True, align="C")
    pdf.ln(6)

    # Reset to normal text colour for body
    pdf.set_text_color(0, 0, 0)

    for section in sections:
        heading = _safe(section.get("heading", ""))
        bullets = section.get("bullets", [])

        if heading:
            pdf.set_font("Helvetica", style="B", size=13)
            pdf.set_fill_color(220, 220, 220)
            pdf.cell(0, 9, heading, new_x="LMARGIN", new_y="NEXT", fill=True)
            pdf.ln(2)

        pdf.set_font("Helvetica", size=11)
        for bullet in bullets:
            # multi_cell wraps long lines automatically.
            # Bullet character is a simple dash to stay ASCII-safe.
            pdf.multi_cell(0, 7, f"  - {_safe(bullet)}", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(3)

    out_path = _timestamped_path("report", ".pdf")
    pdf.output(str(out_path))
    return {"path": str(out_path)}



# Tool: make_pptx

def make_pptx(title: str, slides: list) -> dict:
    """
    Render a PPTX from structured slides using python-pptx.

    Args:
      title   Presentation title, shown on the first (title) slide.
      slides  List of {"heading": str, "bullets": [str]}.
               Each dict becomes one slide with a title and bullet body.

    Returns:
      path  Absolute path to the generated .pptx file.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts
    # Layout 0 = title slide, Layout 1 = title + content
    title_layout   = slide_layouts[0]
    content_layout = slide_layouts[1]

    # --- Title slide ---
    title_slide = prs.slides.add_slide(title_layout)
    # Stretch the title placeholder to full slide width so centering works.
    title_ph = title_slide.shapes.title
    title_ph.left   = Inches(0)
    title_ph.top    = Inches(2.5)
    title_ph.width  = prs.slide_width
    title_ph.text   = title
    for para in title_ph.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
    # Subtitle placeholder (index 1) may not exist on all themes; guard it.
    if len(title_slide.placeholders) > 1:
        sub_ph = title_slide.placeholders[1]
        sub_ph.left  = Inches(0)
        sub_ph.top   = Inches(3.5)
        sub_ph.width = prs.slide_width
        sub_ph.text  = datetime.now().strftime("%Y-%m-%d")
        for para in sub_ph.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER

    # --- Content slides ---
    for slide_data in slides:
        heading = slide_data.get("heading", "")
        bullets = slide_data.get("bullets", [])

        slide = prs.slides.add_slide(content_layout)
        title_shape = slide.shapes.title
        title_shape.left  = Inches(0)
        title_shape.top   = Inches(0.4)
        title_shape.width = prs.slide_width
        title_shape.text  = heading
        for para in title_shape.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER

        # The body placeholder (index 1) holds the bullet text frame.
        body = slide.placeholders[1]
        tf   = body.text_frame
        tf.clear()  # remove the default empty paragraph

        for i, bullet in enumerate(bullets):
            if i == 0:
                # First paragraph already exists after clear(); reuse it.
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text  = bullet
            para.level = 0  # top-level bullet
            para.font.size = Pt(18)

    out_path = _timestamped_path("slides", ".pptx")
    prs.save(str(out_path))
    return {"path": str(out_path)}


# Helper

def _safe(text: str) -> str:
    """Strip characters that PDF/PPTX renderers may reject (null bytes, etc.)."""
    return text.replace("\x00", "").strip() if text else ""


# Entrypoint

if __name__ == "__main__":
    server = MCPServer("generation_server")
    server.register_tool("make_pdf",  make_pdf)
    server.register_tool("make_pptx", make_pptx)
    server.run()
